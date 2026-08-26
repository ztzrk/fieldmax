import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    base_dir = r"c:\Users\Ztzrk\Documents\fieldmax"
    img_dir = os.path.join(base_dir, "obsidian", "Fieldmax", "images")
    output_pptx = os.path.join(base_dir, "obsidian", "Fieldmax", "Presentasi_Seminar_Hasil_FieldMax.pptx")

    # Colors
    BG_CANVAS = RGBColor(248, 245, 238)      # #F8F5EE Warm Ivory Canvas
    TEXT_MAIN = RGBColor(17, 24, 39)         # #111827 Deep Black/Gray
    TEXT_MUTED = RGBColor(75, 85, 99)        # #4B5563 Medium Slate
    CHAPTER_TAG = RGBColor(37, 99, 235)      # #2563EB Chapter Blue
    TABLE_HEADER = RGBColor(37, 99, 235)     # #2563EB Academic Blue
    WHITE = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(209, 213, 219)
    ACCENT_ORANGE = RGBColor(234, 88, 12)    # #EA580C
    ACCENT_GREEN = RGBColor(22, 163, 74)     # #16A34A
    DARK_BOX = RGBColor(17, 24, 39)

    def set_canvas_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_CANVAS
        bg.line.fill.background()
        return bg

    def add_circle_badge(slide, num):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.0), Inches(6.3), Inches(0.65), Inches(0.65))
        badge.fill.solid()
        badge.fill.fore_color.rgb = DARK_BOX
        badge.line.fill.background()
        tf = badge.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = str(num)
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

    def add_header(slide, title, chapter_tag="BAB I • PENDAHULUAN", subtitle=None):
        tb_chap = slide.shapes.add_textbox(Inches(0.9), Inches(0.4), Inches(10.0), Inches(0.35))
        tf_c = tb_chap.text_frame
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = chapter_tag.upper()
        p_c.font.name = "Arial"
        p_c.font.size = Pt(11)
        p_c.font.bold = True
        p_c.font.color.rgb = CHAPTER_TAG

        tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.533), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial Black"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN

        if subtitle:
            ps = tf.add_paragraph()
            ps.space_before = Pt(2)
            ps.text = subtitle
            ps.font.name = "Arial"
            ps.font.size = Pt(14)
            ps.font.bold = True
            ps.font.color.rgb = TEXT_MUTED

    def set_notes(slide, text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text

    # =========================================================================
    # SLIDE 1: COVER
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s1)

    logo_path = os.path.join(img_dir, "image001.png")
    if os.path.exists(logo_path):
        s1.shapes.add_picture(logo_path, Inches(0.9), Inches(0.6), height=Inches(0.75))

    tb_inst = s1.shapes.add_textbox(Inches(1.8), Inches(0.62), Inches(5.0), Inches(0.75))
    tf_i = tb_inst.text_frame
    tf_i.margin_left = tf_i.margin_top = tf_i.margin_right = tf_i.margin_bottom = 0
    p1 = tf_i.paragraphs[0]
    p1.text = "Sistem informasi"
    p1.font.name = "Arial"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    p2 = tf_i.add_paragraph()
    p2.text = "Universitas Hasanuddin"
    p2.font.name = "Arial"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED

    tb_sem = s1.shapes.add_textbox(Inches(9.5), Inches(0.65), Inches(2.9), Inches(0.5))
    tf_s = tb_sem.text_frame
    p_sem = tf_s.paragraphs[0]
    p_sem.alignment = PP_ALIGN.RIGHT
    p_sem.text = "SEMINAR HASIL"
    p_sem.font.name = "Arial Black"
    p_sem.font.size = Pt(14)
    p_sem.font.bold = True
    p_sem.font.color.rgb = TEXT_MAIN

    tb_title = s1.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.533), Inches(2.3))
    tf_t = tb_title.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
    p_t = tf_t.paragraphs[0]
    p_t.alignment = PP_ALIGN.CENTER
    p_t.text = "RANCANG BANGUN SISTEM INFORMASI RESERVASI DAN FASILITAS OLAHRAGA MULTI-TENANT BERBASIS WEB (STUDI KASUS: PLATFORM FIELDMAX)"
    p_t.font.name = "Arial Black"
    p_t.font.size = Pt(25)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_MAIN

    tb_author = s1.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11.533), Inches(1.2))
    tf_a = tb_author.text_frame
    tf_a.word_wrap = True
    tf_a.margin_left = tf_a.margin_top = tf_a.margin_right = tf_a.margin_bottom = 0
    pa1 = tf_a.paragraphs[0]
    pa1.alignment = PP_ALIGN.CENTER
    pa1.text = "OLEH:"
    pa1.font.name = "Arial Black"
    pa1.font.size = Pt(12)
    pa1.font.bold = True
    pa1.font.color.rgb = TEXT_MAIN
    
    pa2 = tf_a.add_paragraph()
    pa2.space_before = Pt(3)
    pa2.alignment = PP_ALIGN.CENTER
    pa2.text = "Aflah Alifu Na Mappatajang Rahman"
    pa2.font.name = "Arial"
    pa2.font.size = Pt(17)
    pa2.font.bold = True
    pa2.font.color.rgb = TEXT_MAIN

    pa3 = tf_a.add_paragraph()
    pa3.alignment = PP_ALIGN.CENTER
    pa3.text = "H071211012"
    pa3.font.name = "Arial"
    pa3.font.size = Pt(14)
    pa3.font.bold = True
    pa3.font.color.rgb = TEXT_MUTED

    tb_dos = s1.shapes.add_textbox(Inches(0.9), Inches(5.8), Inches(5.0), Inches(1.1))
    tf_d = tb_dos.text_frame
    pd1 = tf_d.paragraphs[0]
    pd1.text = "Dosen Pembimbing"
    pd1.font.name = "Arial"
    pd1.font.size = Pt(11)
    pd1.font.bold = True
    pd1.font.color.rgb = TEXT_MAIN
    pd2 = tf_d.add_paragraph()
    pd2.text = "Dr. Eng Supri Bin Hj Amir, S.Si., M.Eng."
    pd2.font.name = "Arial"
    pd2.font.size = Pt(11)
    pd2.font.color.rgb = TEXT_MUTED

    tb_peng = s1.shapes.add_textbox(Inches(6.5), Inches(5.8), Inches(5.0), Inches(1.1))
    tf_p = tb_peng.text_frame
    pp1 = tf_p.paragraphs[0]
    pp1.text = "Penguji"
    pp1.font.name = "Arial"
    pp1.font.size = Pt(11)
    pp1.font.bold = True
    pp1.font.color.rgb = TEXT_MAIN
    pp2 = tf_p.add_paragraph()
    pp2.text = "Edy Saputra Rusdi, S.Si., M.Si.\nDr. Hendra, S.Si., M.Kom."
    pp2.font.name = "Arial"
    pp2.font.size = Pt(11)
    pp2.font.color.rgb = TEXT_MUTED

    add_circle_badge(s1, 1)
    set_notes(s1, "Selamat pagi/siang Dewan Penguji. Saya Aflah Alifu Na Mappatajang Rahman mempresentasikan hasil skripsi FieldMax.")

    # =========================================================================
    # SLIDE 2: BAB I • LATAR BELAKANG (Tren & Fragmentasi Fasilitas)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s2)
    add_header(s2, "Latar Belakang: Tren & Kebutuhan Platform", "BAB I • PENDAHULUAN")

    tb_c = s2.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.0), Inches(4.5))
    tf = tb_c.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Lonjakan kesadaran gaya hidup sehat mendorong peningkatan pesat aktivitas olahraga kelompok (futsal, bulu tangkis, bola basket, dan mini soccer) di perkotaan."
    p1.font.name = "Arial"
    p1.font.size = Pt(22)
    p1.font.color.rgb = TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.space_before = Pt(32)
    p2.text = "Namun, mayoritas sistem yang ada saat ini bersifat terfragmentasi pada satu sarana saja (single-venue), sehingga calon penyewa kesulitan mengeksplorasi perbandingan fasilitas, tarif, dan ketersediaan jadwal secara terpusat."
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_MAIN

    add_circle_badge(s2, 2)
    set_notes(s2, "Subbab 1.1: Kebutuhan platform marketplace terpadu multi-venue.")

    # =========================================================================
    # SLIDE 3: BAB I • LATAR BELAKANG (4 Masalah Krusial)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s3)
    add_header(s3, "Latar Belakang: 4 Masalah Sistem Konvensional", "BAB I • PENDAHULUAN")

    tb_c = s3.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.0), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Proses bisnis konvensional via buku agenda & WhatsApp memiliki 4 masalah krusial (Nadjamuddin 2023, Hafiz 2023):"
    p.font.name = "Arial"
    p.font.size = Pt(19)
    p.font.color.rgb = TEXT_MUTED

    bullets = [
        ("Bentrok Jadwal (Double Booking):", "Ketiadaan kendali konkurensi (concurrency control) memicu jadwal ganda saat jam sibuk."),
        ("Manipulasi Struk Transfer:", "Verifikasi transfer manual lambat dan rentan terhadap pemalsuan bukti transfer fiktif."),
        ("Ketiadaan Informasi Real-Time:", "Calon penyewa harus menunggu konfirmasi berulang kali dari admin pengelola."),
        ("Inefisiensi Rekapitulasi Finansial:", "Mitra pemilik sarana (Renter) kesulitan merekapitulasi omzet dan riwayat sewa secara otomatis.")
    ]

    for title, desc in bullets:
        pb = tf.add_paragraph()
        pb.space_before = Pt(15)
        pb.text = f"•  {title} {desc}"
        pb.font.name = "Arial"
        pb.font.size = Pt(19)
        pb.font.color.rgb = TEXT_MAIN

    add_circle_badge(s3, 3)
    set_notes(s3, "Subbab 1.1: Empat masalah mendasar sistem reservasi konvensional.")

    # =========================================================================
    # SLIDE 4: BAB I • DIAGRAM ANALISIS MASALAH & SOLUSI (MAXIMIZED FULL-WIDTH)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s4)
    add_header(s4, "Analisis Masalah Sistem Manual vs Solusi", "BAB I • PENDAHULUAN")

    prob_img = os.path.join(img_dir, "gambar-analisis-masalah-tight.png")
    if not os.path.exists(prob_img):
        prob_img = os.path.join(img_dir, "gambar-analisis-masalah.png")

    if os.path.exists(prob_img):
        s4.shapes.add_picture(prob_img, Inches(1.2), Inches(1.65), width=Inches(10.933), height=Inches(5.35))

    add_circle_badge(s4, 4)
    set_notes(s4, "Subbab 1.1: Diagram pemetaan masalah sistem manual, dampak operasional, dan solusi terpadu FieldMax.")

    # =========================================================================
    # SLIDE 5: BAB I • RUMUSAN MASALAH
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s5)
    add_header(s5, "Rumusan Masalah", "BAB I • PENDAHULUAN")

    tb_c = s5.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.0), Inches(4.5))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "1. Bagaimana merancang dan membangun sistem informasi reservasi dan fasilitas olahraga multi-tenant berbasis web (FieldMax) yang mampu mengotomatisasi proses bisnis dan menyediakan layanan terpadu?"
    p1.font.name = "Arial"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.space_before = Pt(32)
    p2.text = "2. Bagaimana menguji dan membuktikan keandalan sistem informasi FieldMax menggunakan metode Black Box Testing dalam mengeliminasi terjadinya bentrok jadwal (double booking)?"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MAIN

    add_circle_badge(s5, 5)
    set_notes(s5, "Subbab 1.2: Dua rumusan masalah utama skripsi.")

    # =========================================================================
    # SLIDE 6: BAB I • TUJUAN PENELITIAN
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s6)
    add_header(s6, "Tujuan Penelitian", "BAB I • PENDAHULUAN")

    tb_c = s6.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.0), Inches(4.5))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "1. Merancang dan mengimplementasikan platform FieldMax berbasis Full-Stack TypeScript Monorepo guna mendigitalisasi pengelolaan sarana olahraga bagi User, Renter, dan Admin."
    p1.font.name = "Arial"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.space_before = Pt(32)
    p2.text = "2. Menguji dan membuktikan secara empiris keandalan sistem menggunakan Black Box Testing dalam mencegah dan mengeliminasi terjadinya pemesanan ganda (double booking)."
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MAIN

    add_circle_badge(s6, 6)
    set_notes(s6, "Subbab 1.3: Tujuan penelitian menjawab langsung rumusan masalah.")

    # =========================================================================
    # SLIDE 7: BAB I • BATASAN MASALAH
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s7)
    add_header(s7, "Batasan Masalah", "BAB I • PENDAHULUAN")

    tb_c = s7.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.0), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    b_items = [
        "1. Arsitektur Web Responsif: Next.js 16 (App Router) pada Front End dan Express.js 5 pada Back End berbasis TypeScript.",
        "2. Basis Data & Transaksi: PostgreSQL 16 dengan Prisma ORM, serta pembayaran digital Midtrans Snap (QRIS, VA, e-Wallet).",
        "3. Layanan Pendukung: ImageKit CDN (media foto), Nodemailer SMTP (token email), dan otentikasi berbasis sesi di database.",
        "4. Ruang Lingkup Uji: Fungsionalitas Black Box Testing untuk verifikasi proses bisnis dan eliminasi double booking."
    ]

    for idx, item in enumerate(b_items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if idx > 0:
            p.space_before = Pt(20)
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s7, 7)
    set_notes(s7, "Subbab 1.4: Batasan teknologi dan ruang lingkup pengujian Black Box.")

    # =========================================================================
    # SLIDE 8: BAB I • MANFAAT PENELITIAN
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s8)
    add_header(s8, "Manfaat Penelitian", "BAB I • PENDAHULUAN")

    tb_c = s8.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.0), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    m_items = [
        ("1. Bagi Pelanggan (User):", "Kepastian ketersediaan jadwal real-time, eksplorasi multi-cabor, dan transaksi digital mandiri 24/7."),
        ("2. Bagi Pemilik Sarana (Renter):", "Otonomi manajemen jadwal & tarif per jam, zero double booking, serta analitik omzet otomatis."),
        ("3. Bagi Administrator (Admin):", "Kontrol terpusat moderasi legalitas venue/lapangan serta penanganan tiket pengaduan."),
        ("4. Bagi Teoretis & Akademik:", "Rujukan penerapan DSR dan Waterfall SDLC pada sistem informasi multi-tenant modern.")
    ]

    for idx, (title, desc) in enumerate(m_items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if idx > 0:
            p.space_before = Pt(18)
        p.text = f"{title} {desc}"
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s8, 8)
    set_notes(s8, "Subbab 1.5: Manfaat penelitian bagi empat pihak terkait.")

    # =========================================================================
    # SLIDE 9: BAB II • WAKTU DAN LOKASI PENELITIAN
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s9)
    add_header(s9, "Waktu dan Lokasi Penelitian", "BAB II • METODE PENELITIAN")

    tb_c = s9.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(0.8))
    tf = tb_c.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Penelitian dilaksanakan pada bulan Juli 2025 sampai November 2025 di Kota Makassar, Sulawesi Selatan (Tabel 4 Skripsi)."
    p.font.name = "Arial"
    p.font.size = Pt(17)
    p.font.color.rgb = TEXT_MAIN

    table_shape = s9.shapes.add_table(7, 7, Inches(0.9), Inches(2.6), Inches(11.5), Inches(4.0))
    table = table_shape.table
    table.columns[0].width = Inches(0.7)
    table.columns[1].width = Inches(4.8)
    for c in range(2, 7):
        table.columns[c].width = Inches(1.2)

    gantt_headers = ["No", "Tahapan Penelitian", "Juli", "Agustus", "September", "Oktober", "November"]
    for i, h in enumerate(gantt_headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_ORANGE
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if i != 1 else PP_ALIGN.LEFT
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE

    gantt_rows = [
        ("1", "Studi Literatur", [True, True, False, False, False]),
        ("2", "Analisis Kebutuhan (Requirements)", [False, True, False, False, False]),
        ("3", "Desain Sistem (Design)", [False, True, True, False, False]),
        ("4", "Implementasi Sistem (Implementation)", [False, False, True, True, False]),
        ("5", "Pengujian Sistem (Testing)", [False, False, False, True, True]),
        ("6", "Pemeliharaan & Dokumentasi (Maintenance)", [False, False, False, False, True]),
    ]

    for r_idx, (num, task, filled) in enumerate(gantt_rows, start=1):
        table.cell(r_idx, 0).text = num
        table.cell(r_idx, 1).text = task
        for col_idx, is_filled in enumerate(filled, start=2):
            cell = table.cell(r_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT_ORANGE if is_filled else WHITE
            if is_filled:
                cell.text_frame.paragraphs[0].text = "✓"
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = WHITE

    add_circle_badge(s9, 9)
    set_notes(s9, "Subbab 2.1 & Tabel 4: Jadwal pelaksanaan penelitian Juli - November 2025.")

    # =========================================================================
    # SLIDE 10: BAB II • DESIGN SCIENCE RESEARCH (MAXIMIZED FULL-WIDTH)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s10)
    add_header(s10, "Kerangka Design Science Research (DSR)", "BAB II • METODE PENELITIAN")

    dsr_img = os.path.join(img_dir, "gambar-dsr-fieldmax-tight.png")
    if not os.path.exists(dsr_img):
        dsr_img = os.path.join(img_dir, "gambar-dsr-fieldmax.png")

    if os.path.exists(dsr_img):
        s10.shapes.add_picture(dsr_img, Inches(1.2), Inches(1.65), width=Inches(10.933), height=Inches(5.35))

    add_circle_badge(s10, 10)
    set_notes(s10, "Subbab 2.2: Kerangka DSR Hevner et al. (2004) menghubungkan Environment, IS Research, dan Knowledge Base.")

    # =========================================================================
    # SLIDE 11: BAB II • METODE PENGUMPULAN DATA
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s11)
    add_header(s11, "Metode Pengumpulan Data: Studi Literatur", "BAB II • METODE PENELITIAN")

    tb_c = s11.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.0), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Pengumpulan data dilakukan secara komprehensif melalui Studi Literatur (Literature Review):"
    p.font.name = "Arial"
    p.font.size = Pt(19)
    p.font.color.rgb = TEXT_MUTED

    l_items = [
        ("1. Jurnal Ilmiah Terakreditasi:", "Mengkaji isu konkurensi, eliminasi double booking, arsitektur multi-tenant, dan integrasi payment gateway."),
        ("2. Buku Teks RPL & Sistem Informasi:", "Teori pemodelan UML (Pressman, 2010), Waterfall SDLC (Pfleeger & Atlee, 2010), dan kerangka DSR (Hevner et al., 2004)."),
        ("3. Dokumentasi Teknis Resmi:", "Spesifikasi Next.js 16 App Router, Express.js 5, PostgreSQL 16, Prisma ORM, Midtrans API, dan ImageKit SDK.")
    ]

    for title, desc in l_items:
        pb = tf.add_paragraph()
        pb.space_before = Pt(18)
        pb.text = f"{title} {desc}"
        pb.font.name = "Arial"
        pb.font.size = Pt(19)
        pb.font.color.rgb = TEXT_MAIN

    add_circle_badge(s11, 11)
    set_notes(s11, "Subbab 2.3: Tiga domain studi literatur ilmiah.")

    # =========================================================================
    # SLIDE 12: BAB II • WATERFALL SDLC
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s12)
    add_header(s12, "Metode Pengembangan: Waterfall SDLC", "BAB II • METODE PENELITIAN")

    tb_c = s12.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.0), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    wf_steps = [
        ("1. Requirements:", "Analisis 10 kendala sistem manual dan elisitasi 22 Use Case fungsional."),
        ("2. Design:", "Perancangan Use Case, Activity Diagram, antarmuka Figma, dan 16 relasi ERD."),
        ("3. Implementation:", "Pengkodean Next.js 16, Express.js 5, Prisma ORM, & PostgreSQL 16."),
        ("4. Testing:", "Pengujian Black Box 76 kasus uji & 7 skenario matematis bentrok jadwal."),
        ("5. Maintenance:", "Optimasi performa kueri basis data dan pemeliharaan server berkala.")
    ]

    for idx, (title, desc) in enumerate(wf_steps):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if idx > 0:
            p.space_before = Pt(16)
        p.text = f"{title} {desc}"
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s12, 12)
    set_notes(s12, "Subbab 2.4: 5 fase Waterfall SDLC Pfleeger & Atlee (2010).")

    # =========================================================================
    # SLIDE 13: BAB II • TAHAPAN PENELITIAN (FLOWCHART + MILESTONES)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s13)
    add_header(s13, "Tahapan Penelitian (Alur Sistematis)", "BAB II • METODE PENELITIAN")

    vec_img = os.path.join(img_dir, "gambar-4-alur-penelitian-vector.png")
    if os.path.exists(vec_img):
        s13.shapes.add_picture(vec_img, Inches(1.3), Inches(1.65), height=Inches(5.4))

    tb_flow = s13.shapes.add_textbox(Inches(4.1), Inches(1.7), Inches(7.8), Inches(5.2))
    tf_f = tb_flow.text_frame
    tf_f.word_wrap = True

    steps_milestone = [
        ("1. Mulai & Studi Literatur:", "Identifikasi masalah riil reservasi & landasan teori DSR Hevner."),
        ("2. Requirements (Analisis):", "Elisitasi kebutuhan fungsional 3 peran (User, Renter, Admin)."),
        ("3. Design (Perancangan):", "Pemodelan 22 Use Case, Activity Diagram, Figma UI, & 16 tabel ERD."),
        ("4. Implementation (Pengkodean):", "Next.js 16 App Router, Express.js 5, Prisma ORM, & PostgreSQL."),
        ("5. Testing (Evaluasi Black Box):", "Uji 76 skenario Black Box & 7 formula irisan waktu bentrok jadwal."),
        ("6. Hasil Penelitian & Selesai:", "Dokumentasi naskah komprehensif, publikasi, dan deployment.")
    ]

    for idx, (title, desc) in enumerate(steps_milestone):
        p = tf_f.paragraphs[0] if idx == 0 else tf_f.add_paragraph()
        if idx > 0:
            p.space_before = Pt(13)
        p.text = f"{title} {desc}"
        p.font.name = "Arial"
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s13, 13)
    set_notes(s13, "Subbab 2.5 & Gambar 6: Alur tahapan penelitian terintegrasi DSR dan Waterfall SDLC.")

    # =========================================================================
    # SLIDE 14: BAB II • PERANCANGAN SISTEM (22 USE CASES)
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s14)
    add_header(s14, "Perancangan Sistem: 22 Use Cases Terpadu", "BAB II • PERANCANGAN SISTEM")

    uc_data = [
        ("1. User / Pelanggan (8)", TABLE_HEADER, [
            "1. Daftar dan Login Akun",
            "2. Cari dan Filter Lapangan",
            "3. Lihat Detail Venue & Lapangan",
            "4. Reservasi Slot Jadwal Real-Time",
            "5. Bayar via Midtrans Snap",
            "6. Lihat Riwayat Pemesanan",
            "7. Beri Rating Ulasan Lapangan",
            "8. Buat Tiket Pengaduan Kendala"
        ]),
        ("2. Mitra Renter (7)", ACCENT_GREEN, [
            "1. Daftar dan Login Mitra Bisnis",
            "2. Kelola Profil Venue & Jam Operasional",
            "3. Kelola Lapangan & Tarif per Jam",
            "4. Ajukan Venue & Lapangan ke Admin",
            "5. Kelola & Pantau Pemesanan Masuk",
            "6. Lihat Analitik Grafik Pendapatan",
            "7. Kirim Tiket Pengaduan ke Admin",
            ""
        ]),
        ("3. Administrator (7)", ACCENT_ORANGE, [
            "1. Login Panel Administrasi",
            "2. Lihat Dashboard Statistik Ekosistem",
            "3. Kelola & Verifikasi Akun Pengguna",
            "4. Kelola Master Data Sport Types",
            "5. Moderasi Venue dan Unit Lapangan",
            "6. Pantau Pemesanan & Pembayaran",
            "7. Kelola & Resolusi Tiket Pengaduan",
            ""
        ])
    ]

    for idx, (header_title, header_col, items) in enumerate(uc_data):
        x = Inches(0.9 + idx * 3.9)
        card = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(3.65), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)

        tb = s14.shapes.add_textbox(x + Inches(0.18), Inches(2.15), Inches(3.3), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_h = tf.paragraphs[0]
        p_h.text = header_title
        p_h.font.name = "Arial Black"
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = header_col

        for item in items:
            if item:
                pi = tf.add_paragraph()
                pi.space_before = Pt(8)
                pi.text = item
                pi.font.name = "Arial"
                pi.font.size = Pt(12)
                pi.font.color.rgb = TEXT_MAIN

    add_circle_badge(s14, 14)
    set_notes(s14, "Subbab 2.7: Rincian lengkap 22 Use Case yang terdistribusi ke 3 aktor utama.")

    # =========================================================================
    # SLIDE 15: BAB III • IMPLEMENTASI FRONT-END & ALASAN PEMILIHAN
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s15)
    add_header(s15, "Implementasi Front-End & Alasan Pemilihan", "BAB III • HASIL DAN PEMBAHASAN")

    tb_c = s15.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p_badge = tf.paragraphs[0]
    p_badge.text = "STACK: Next.js 16 (App Router)   •   React 19   •   Tailwind CSS v4   •   TanStack Query"
    p_badge.font.name = "Arial Black"
    p_badge.font.size = Pt(16)
    p_badge.font.bold = True
    p_badge.font.color.rgb = TABLE_HEADER

    fe_reasons = [
        ("1. Server-Side Rendering (SSR) & App Router:", "Mempercepat waktu render awal (initial load) katalog venue, pemisahan rute dinamis modular, dan performa tinggi ramah SEO."),
        ("2. Tailwind CSS v4 & shadcn/ui:", "Menjamin konsistensi visual modern, komponen antarmuka yang aksesibel, serta tata letak responsif di peramban mobile dan desktop."),
        ("3. TanStack React Query:", "Manajemen state asinkron dan caching otomatis ketersediaan jadwal slot sewa, menghemat beban lalu lintas data ke peladen.")
    ]

    for title, desc in fe_reasons:
        p = tf.add_paragraph()
        p.space_before = Pt(18)
        p.text = f"{title} {desc}"
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s15, 15)
    set_notes(s15, "Subbab 3.1 & 1.6.4: Arsitektur Front-End dan alasan pemilihannya.")

    # =========================================================================
    # SLIDE 16: BAB III • IMPLEMENTASI BACK-END & ALASAN PEMILIHAN
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s16)
    add_header(s16, "Implementasi Back-End & Alasan Pemilihan", "BAB III • HASIL DAN PEMBAHASAN")

    tb_c = s16.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p_badge = tf.paragraphs[0]
    p_badge.text = "STACK: Express.js 5   •   Three-Tier Layered   •   Zod Validation   •   Session Auth"
    p_badge.font.name = "Arial Black"
    p_badge.font.size = Pt(16)
    p_badge.font.bold = True
    p_badge.font.color.rgb = TABLE_HEADER

    be_reasons = [
        ("1. Pola Three-Tier (Routes -> Controllers -> Services):", "Memisahkan logika bisnis dari protokol HTTP dan akses data sehingga sistem modular, terstruktur, dan mudah dirawat."),
        ("2. Full-Stack TypeScript & Zod:", "Memberikan jaminan keamanan tipe data (End-to-End Type Safety) antara front-end dan back-end guna mencegah galat tipe data API."),
        ("3. Session-Based Auth (HttpOnly Cookie):", "Menyimpan status otentikasi aman di basis data, melindungi sesi pengguna dari serangan pencurian token (XSS).")
    ]

    for title, desc in be_reasons:
        p = tf.add_paragraph()
        p.space_before = Pt(18)
        p.text = f"{title} {desc}"
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s16, 16)
    set_notes(s16, "Subbab 3.1 & 1.6.4: Arsitektur Back-End dan alasan pemilihannya.")

    # =========================================================================
    # SLIDE 17: BAB III • IMPLEMENTASI BASIS DATA & ALASAN PEMILIHAN
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s17)
    add_header(s17, "Implementasi Basis Data & Alasan Pemilihan", "BAB III • HASIL DAN PEMBAHASAN")

    tb_c = s17.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p_badge = tf.paragraphs[0]
    p_badge.text = "STACK: PostgreSQL 16   •   Prisma ORM   •   MVCC & ACID   •   UUID Primary Key"
    p_badge.font.name = "Arial Black"
    p_badge.font.size = Pt(16)
    p_badge.font.bold = True
    p_badge.font.color.rgb = TABLE_HEADER

    db_reasons = [
        ("1. PostgreSQL MVCC & Transaksi ACID:", "Menangani konkurensi reservasi secara atomik saat banyak pengguna memesan slot waktu secara serentak tanpa deadlock data."),
        ("2. Prisma ORM (Schema Migration & Client):", "Menyediakan akses data type-safe, migrasi deklaratif otomatis, serta proteksi terintegrasi dari SQL Injection."),
        ("3. UUID (Universally Unique Identifier):", "Mengamankan entitas data dari serangan enumerasi ID berurutan dan menjamin keunikan kunci primer lintas tabel.")
    ]

    for title, desc in db_reasons:
        p = tf.add_paragraph()
        p.space_before = Pt(18)
        p.text = f"{title} {desc}"
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s17, 17)
    set_notes(s17, "Subbab 3.1, 3.2, & 1.6.4: Basis Data dan alasan pemilihannya.")

    # =========================================================================
    # SLIDE 18: BAB III • INTEGRASI LAYANAN PIHAK KETIGA & ALASAN
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s18)
    add_header(s18, "Integrasi Layanan Pihak Ketiga & Alasan", "BAB III • HASIL DAN PEMBAHASAN")

    tb_c = s18.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p_badge = tf.paragraphs[0]
    p_badge.text = "STACK: Midtrans Snap API   •   ImageKit CDN   •   SMTP Nodemailer"
    p_badge.font.name = "Arial Black"
    p_badge.font.size = Pt(16)
    p_badge.font.bold = True
    p_badge.font.color.rgb = TABLE_HEADER

    tp_reasons = [
        ("1. Midtrans Snap API (Payment Gateway):", "Otomasi verifikasi pembayaran digital Indonesia (QRIS, VA, e-Wallet) via Webhook Callback tanpa verifikasi manual bukti struk."),
        ("2. ImageKit CDN (Media Cloud Storage):", "Kompresi dan pengoptimalan resolusi foto fasilitas olahraga secara real-time, mempercepat pemuatan halaman web."),
        ("3. SMTP Nodemailer (Email Notification):", "Otomasi pengiriman kode OTP verifikasi akun 6 digit dan tautan pemulihan kata sandi (reset password) secara instan.")
    ]

    for title, desc in tp_reasons:
        p = tf.add_paragraph()
        p.space_before = Pt(18)
        p.text = f"{title} {desc}"
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s18, 18)
    set_notes(s18, "Subbab 3.1 & 1.6.4: Layanan pihak ketiga dan alasan pemilihannya.")

    # =========================================================================
    # SLIDE 19: BAB III • BASIS DATA - ERD 16 TABEL (MAXIMIZED FULL-WIDTH)
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s19)
    add_header(s19, "Implementasi Basis Data: 16 Tabel Relasional", "BAB III • HASIL DAN PEMBAHASAN")

    erd_img = os.path.join(img_dir, "gambar-relasi-tabel.png")
    if os.path.exists(erd_img):
        s19.shapes.add_picture(erd_img, Inches(1.4), Inches(1.65), width=Inches(10.533), height=Inches(5.35))

    add_circle_badge(s19, 19)
    set_notes(s19, "Subbab 3.2.1 & Gambar 9a: ERD 16 tabel relasional.")

    # =========================================================================
    # SLIDE 20: BAB III • DAFTAR NILAI ENUM (Tabel 5 Skripsi)
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s20)
    add_header(s20, "Implementasi Basis Data: Daftar Enum (Tabel 5)", "BAB III • HASIL DAN PEMBAHASAN")

    enums = [
        ("UserRole {", "  USER\n  RENTER\n  ADMIN\n}"),
        ("BookingStatus {", "  PENDING\n  CONFIRMED\n  CANCELLED\n  COMPLETED\n}"),
        ("PaymentStatus {", "  PENDING\n  PAID\n  EXPIRED\n  FAILED\n}"),
        ("VerificationStatus {", "  DRAFT\n  PENDING\n  APPROVED\n  REJECTED\n}"),
        ("ReportStatus {", "  PENDING\n  RESOLVED\n}"),
        ("ReportCategory {", "  SCAM\n  TECHNICAL\n  PAYMENT\n  OTHER\n}")
    ]

    for idx, (title, body) in enumerate(enums):
        col_idx = idx % 3
        row_idx = idx // 3
        x = Inches(0.9 + col_idx * 3.9)
        y = Inches(2.2 + row_idx * 2.1)

        card = s20.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.6), Inches(1.9))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BOX
        card.line.fill.background()

        tb = s20.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(3.2), Inches(1.6))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Consolas"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(96, 165, 250)

        pb = tf.add_paragraph()
        pb.space_before = Pt(4)
        pb.text = body
        pb.font.name = "Consolas"
        pb.font.size = Pt(11)
        pb.font.color.rgb = WHITE

    add_circle_badge(s20, 20)
    set_notes(s20, "Subbab 3.2.2 & Tabel 5: Enam tipe Enum Prisma ORM.")

    # =========================================================================
    # SLIDE 21: BAB III • STRUKTUR TABEL USERS & SESSIONS (Tabel 6)
    # =========================================================================
    s21 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s21)
    add_header(s21, "Struktur Tabel: users & sessions (Tabel 6)", "BAB III • HASIL DAN PEMBAHASAN")

    t_shape = s21.shapes.add_table(6, 4, Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.4))
    table = t_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(3.8)
    table.columns[3].width = Inches(2.7)

    headers = ["Nama Field", "Tipe Field", "Keterangan", "Default"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE

    t_data = [
        ("id", "String (UUID)", "Primary Key", "uuid()"),
        ("full_name", "String", "Nama lengkap pengguna", "No Default"),
        ("email", "String", "Email unik untuk login", "No Default"),
        ("password", "String", "Hash password akun (bcrypt)", "No Default"),
        ("role", "UserRole", "Hak akses (USER/RENTER/ADMIN)", "USER"),
    ]

    for r_idx, (f, t, k, d) in enumerate(t_data, start=1):
        for c_idx, val in enumerate([f, t, k, d]):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else RGBColor(243, 244, 246)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s21, 21)
    set_notes(s21, "Subbab 3.2.2 & Tabel 6: Struktur kamus data tabel users.")

    # =========================================================================
    # SLIDE 22: BAB III • STRUKTUR TABEL VENUES & BOOKINGS (Tabel 11 & 16)
    # =========================================================================
    s22 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s22)
    add_header(s22, "Struktur Tabel: venues & bookings (Tabel 11 & 16)", "BAB III • HASIL DAN PEMBAHASAN")

    t_shape = s22.shapes.add_table(6, 4, Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.4))
    table = t_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(3.8)
    table.columns[3].width = Inches(2.7)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE

    vf_data = [
        ("venues.id", "String (UUID)", "Primary Key", "uuid()"),
        ("venues.renter_id", "String", "Foreign Key ke users.id", "No Default"),
        ("venues.status", "VerificationStatus", "Status persetujuan admin", "DRAFT"),
        ("bookings.id", "String (UUID)", "Primary Key", "uuid()"),
        ("bookings.status", "BookingStatus", "Status transaksi reservasi", "PENDING"),
    ]

    for r_idx, (f, t, k, d) in enumerate(vf_data, start=1):
        for c_idx, val in enumerate([f, t, k, d]):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else RGBColor(243, 244, 246)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s22, 22)
    set_notes(s22, "Subbab 3.2.2 & Tabel 11 & 16: Struktur tabel venues dan bookings.")

    # =========================================================================
    # SLIDE 23: BAB III • LOGIKA PENCEGAHAN BENTROK JADWAL
    # =========================================================================
    s23 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s23)
    add_header(s23, "Logika Pencegahan Bentrok Jadwal", "BAB III • HASIL DAN PEMBAHASAN")

    box = s23.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BOX
    box.line.color.rgb = TABLE_HEADER
    box.line.width = Pt(2)
    tf_b = box.text_frame
    p = tf_b.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "Kondisi Bentrok = (Start_req < End_db) ∧ (End_req > Start_db)"
    p.font.name = "Consolas"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(96, 165, 250)

    tb_c = s23.shapes.add_textbox(Inches(0.9), Inches(3.7), Inches(11.5), Inches(3.2))
    tf = tb_c.text_frame
    tf.word_wrap = True

    pts = [
        "1. Filter Status Kueri: Memeriksa reservasi aktif berstatus PENDING dan CONFIRMED.",
        "2. Presisi Batas Jam: Mengizinkan jam berdampingan (Adjacent Boundary, misal 08.00–10.00 dan 10.00–12.00).",
        "3. Two-Phase Status Locking: Slot terkunci 15 menit saat checkout dan otomatis lepas (Auto-Release) bila kedaluwarsa."
    ]

    for idx, pt in enumerate(pts):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if idx > 0:
            p.space_before = Pt(16)
        p.text = pt
        p.font.name = "Arial"
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s23, 23)
    set_notes(s23, "Subbab 3.5.1: Formulasi matematis kueri relasional irisan waktu.")

    # =========================================================================
    # SLIDE 24: BAB III • HASIL UJI 7 SKENARIO BENTROK
    # =========================================================================
    s24 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s24)
    add_header(s24, "Hasil Evaluasi 7 Skenario Bentrok (Tabel 22)", "BAB III • HASIL DAN PEMBAHASAN")

    table_shape = s24.shapes.add_table(8, 4, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(3.4)
    table.columns[2].width = Inches(4.8)
    table.columns[3].width = Inches(2.7)

    headers22 = ["No", "Skenario Uji Bentrok (Case)", "Kondisi DB vs Permintaan Baru", "Hasil Evaluasi"]
    for i, h in enumerate(headers22):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if i in [0, 3] else PP_ALIGN.LEFT
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE

    cases_data = [
        ("1", "Exact Slot Match (Jam Sama)", "Eksis: 08.00–10.00 (CONFIRMED) | Baru: 08.00–10.00", "100% Valid (Ditolak)"),
        ("2", "Start-Time Overlap (Tumpang Awal)", "Eksis: 08.00–10.00 (CONFIRMED) | Baru: 07.00–09.00", "100% Valid (Ditolak)"),
        ("3", "End-Time Overlap (Tumpang Akhir)", "Eksis: 08.00–10.00 (CONFIRMED) | Baru: 09.00–11.00", "100% Valid (Ditolak)"),
        ("4", "Enclosing Overlap (Tumpang Tengah)", "Eksis: 08.00–12.00 (PENDING) | Baru: 09.00–11.00", "100% Valid (Ditolak)"),
        ("5", "Adjacent Boundary (Berdampingan)", "Eksis: 08.00–10.00 (CONFIRMED) | Baru: 10.00–12.00", "100% Valid (Diterima)"),
        ("6", "Simultaneous Booking (Serentak)", "Dua user memesan slot 14.00–16.00 bersamaan", "100% Valid (1 Terima, 1 Tolak)"),
        ("7", "Auto-Release on Expire/Cancel", "Slot 16.00–18.00 PENDING kedaluwarsa", "100% Valid (Slot Terbuka)"),
    ]

    for r_idx, (num, name, cond, res) in enumerate(cases_data, start=1):
        for c_idx, val in enumerate([num, name, cond, res]):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else RGBColor(243, 244, 246)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            if c_idx == 0:
                p.alignment = PP_ALIGN.CENTER
            elif c_idx == 3:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True
                p.font.color.rgb = ACCENT_GREEN
            else:
                p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s24, 24)
    set_notes(s24, "Subbab 3.5.1 & Tabel 22: Seluruh 7 skenario uji bentrok jadwal terbukti 100% valid.")

    # =========================================================================
    # SLIDE 25: BAB III • HASIL BLACK BOX TESTING (6 MODUL SISTEM)
    # =========================================================================
    s25 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s25)
    add_header(s25, "Hasil Pengujian Black Box (6 Modul Sistem)", "BAB III • HASIL DAN PEMBAHASAN")

    table_shape = s25.shapes.add_table(7, 4, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(4.3)
    table.columns[2].width = Inches(4.1)
    table.columns[3].width = Inches(2.5)

    bb_headers = ["No", "Modul Pengujian (Tabel Skripsi)", "Cakupan Pengujian Fungsional", "Hasil Evaluasi"]
    for i, h in enumerate(bb_headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if i in [0, 3] else PP_ALIGN.LEFT
        p.text = h
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE

    bb_data = [
        ("1", "Halaman Utama & Pencarian (Tabel 22)", "Filter cabor, geolokasi venue, detail tarif & ketersediaan", "100% Valid (Berhasil)"),
        ("2", "Otentikasi & Manajemen Akun (Tabel 23)", "Registrasi, OTP email, login session, proteksi role RBAC", "100% Valid (Berhasil)"),
        ("3", "Reservasi & Pembayaran Midtrans (Tabel 24)", "Booking slot waktu, invoice Snap Midtrans, webhook otomatis", "100% Valid (Berhasil)"),
        ("4", "Ulasan Lapangan & Pengaduan (Tabel 25)", "Rating bintang, ulasan pengguna, tiket komplain & respon", "100% Valid (Berhasil)"),
        ("5", "Pengelolaan Sarana Mitra/Renter (Tabel 26)", "CRUD venue, jam operasional, tarif lapangan, analitik omzet", "100% Valid (Berhasil)"),
        ("6", "Panel Moderasi Administrator (Tabel 27)", "Verifikasi legalitas venue, moderasi lapangan, master sport", "100% Valid (Berhasil)"),
    ]

    for r_idx, (num, mod, cases, res) in enumerate(bb_data, start=1):
        for c_idx, val in enumerate([num, mod, cases, res]):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else RGBColor(243, 244, 246)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.text = val
            p.font.name = "Arial"
            p.font.size = Pt(10.5)
            if c_idx == 0:
                p.alignment = PP_ALIGN.CENTER
            elif c_idx == 3:
                p.alignment = PP_ALIGN.CENTER
                p.font.bold = True
                p.font.color.rgb = ACCENT_GREEN
            else:
                p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s25, 25)
    set_notes(s25, "Subbab 3.5, Subbab 4.1 & Tabel 22-27: Pengujian fungsionalitas Black Box Testing pada 6 modul utama sistem terbukti 100% valid.")

    # =========================================================================
    # SLIDE 26: BAB III • UI SHOWCASE GUEST & USER
    # =========================================================================
    s26 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s26)
    add_header(s26, "User Interface: Guest & User", "BAB III • HASIL DAN PEMBAHASAN")

    ui_user = [
        ("halaman-utama.png", "1. Halaman Utama (Pencarian)"),
        ("halaman-detail-lapangan.png", "2. Detail Lapangan & Slot Real-Time"),
        ("halaman-riwayat-booking.png", "3. Riwayat Sewa & Invoice")
    ]

    for idx, (img_name, cap) in enumerate(ui_user):
        x = Inches(0.9 + idx * 3.9)
        img_p = os.path.join(img_dir, img_name)
        if os.path.exists(img_p):
            s26.shapes.add_picture(img_p, x, Inches(1.8), width=Inches(3.6), height=Inches(4.5))

        tb = s26.shapes.add_textbox(x, Inches(6.4), Inches(3.6), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = cap
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s26, 26)
    set_notes(s26, "Subbab 3.4: Showcase antarmuka pengguna sisi User.")

    # =========================================================================
    # SLIDE 27: BAB III • UI SHOWCASE RENTER & ADMIN
    # =========================================================================
    s27 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s27)
    add_header(s27, "User Interface: Renter & Administrator", "BAB III • HASIL DAN PEMBAHASAN")

    ui_ra = [
        ("halaman-dashboard-renter.png", "1. Dashboard Mitra (Renter)"),
        ("halaman-pendapatan-renter.png", "2. Analitik Omzet & Pendapatan"),
        ("halaman-moderasi-venue-admin.png", "3. Panel Moderasi Legalitas (Admin)")
    ]

    for idx, (img_name, cap) in enumerate(ui_ra):
        x = Inches(0.9 + idx * 3.9)
        img_p = os.path.join(img_dir, img_name)
        if os.path.exists(img_p):
            s27.shapes.add_picture(img_p, x, Inches(1.8), width=Inches(3.6), height=Inches(4.5))

        tb = s27.shapes.add_textbox(x, Inches(6.4), Inches(3.6), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = cap
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s27, 27)
    set_notes(s27, "Subbab 3.4: Showcase antarmuka mitra dan administrator.")

    # =========================================================================
    # SLIDE 28: BAB IV • KESIMPULAN PENELITIAN
    # =========================================================================
    s28 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s28)
    add_header(s28, "Kesimpulan Penelitian", "BAB IV • KESIMPULAN & SARAN")

    tb_c = s28.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.0), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "1. Rancang Bangun Sistem Terpadu (Build): Platform FieldMax berhasil dibangun dengan Full-Stack TypeScript Monorepo (Next.js 16, Express.js 5, Prisma ORM, PostgreSQL 16, dan Midtrans Snap) yang mengintegrasikan 22 use cases untuk 3 peran pengguna secara efisien."
    p1.font.name = "Arial"
    p1.font.size = Pt(20)
    p1.font.color.rgb = TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.space_before = Pt(26)
    p2.text = "2. Evaluasi Eliminasi Double Booking (Evaluate): Pengujian fungsionalitas Black Box Testing pada 6 modul utama (Tabel 22–27) dan validasi empiris 7 skenario irisan waktu terbukti 100% valid mengeliminasi potensi jadwal ganda (double booking) serta membebaskan kembali slot sewa yang kedaluwarsa secara otomatis."
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_MAIN

    add_circle_badge(s28, 28)
    set_notes(s28, "Subbab 4.1: Kesimpulan menjawab rumusan masalah 1 dan 2.")

    # =========================================================================
    # SLIDE 29: BAB IV • SARAN & REKOMENDASI
    # =========================================================================
    s29 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s29)
    add_header(s29, "Saran & Rekomendasi Selanjutnya", "BAB IV • KESIMPULAN & SARAN")

    tb_c = s29.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.0), Inches(4.8))
    tf = tb_c.text_frame
    tf.word_wrap = True

    s_items = [
        "1. Aplikasi Mobile (Native / PWA): Pengembangan aplikasi bergerak yang dilengkapi push notification jadwal bermain.",
        "2. Peta Interaktif GIS (Maps API): Navigasi visual dan pencarian sarana terdekat berbasis geolokasi GPS pengguna.",
        "3. Komunitas & Sparring Matchmaking: Fitur pencarian lawan tanding (sparring) dan pembagian tagihan sewa (split bill).",
        "4. Gateway WhatsApp Business API: Notifikasi kode booking dan bukti bayar digital instan via WhatsApp.",
        "5. Pengujian Beban (Stress Testing): Uji ketahanan database saat lonjakan transaksi masif pada jam-jam sibuk (peak hours)."
    ]

    for idx, item in enumerate(s_items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if idx > 0:
            p.space_before = Pt(14)
        p.text = item
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_MAIN

    add_circle_badge(s29, 29)
    set_notes(s29, "Subbab 4.2: Lima saran rekomendasi pengembangan masa depan.")

    # =========================================================================
    # SLIDE 30: PENUTUP / THANK YOU
    # =========================================================================
    s30 = prs.slides.add_slide(blank_layout)
    set_canvas_bg(s30)

    tb_sec = s30.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.0), Inches(1.5))
    tf = tb_sec.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank you"
    p.font.name = "Arial Black"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN

    tb_sub = s30.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(10.0), Inches(1.2))
    tf_s = tb_sub.text_frame
    p_sub = tf_s.paragraphs[0]
    p_sub.text = "Sesi Tanya Jawab, Kritik & Saran Dewan Penguji"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(22)
    p_sub.font.bold = True
    p_sub.font.color.rgb = CHAPTER_TAG

    p_cred = tf_s.add_paragraph()
    p_cred.space_before = Pt(12)
    p_cred.text = "Aflah Alifu Na Mappatajang Rahman • H071211012 • Sistem Informasi FMIPA Unhas 2025"
    p_cred.font.name = "Arial"
    p_cred.font.size = Pt(15)
    p_cred.font.color.rgb = TEXT_MUTED

    add_circle_badge(s30, 30)
    set_notes(s30, "Sekian presentasi dari saya. Terima kasih kepada Dewan Penguji. Sesi tanya jawab dipersilakan.")

    prs.save(output_pptx)
    print(f"SUCCESS: 30-Slide presentation with maximized diagrams generated at:\n{output_pptx}")
    return output_pptx

if __name__ == "__main__":
    build_presentation()
